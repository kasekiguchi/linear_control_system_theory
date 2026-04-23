#!/usr/bin/env python3
"""Generalized chapter -> pptx converter.

Usage: python3 build_chap.py N
  N = chapter number (1,2,3,4)

Uses chap3.pptx as the reference-doc so the user's slide master is reused.
"""
import re, subprocess, sys, pathlib, pypandoc

N = int(sys.argv[1])
ROOT = pathlib.Path(__file__).parent
BASE = ROOT.parent
SRC  = BASE / f"chap{N}.tex"
OUT  = ROOT / f"chap{N}.pptx"
WRAP = ROOT / f"_chap{N}_wrapped.tex"
MD   = ROOT / f"chap{N}.md"
REF  = ROOT / "chap3.pptx"   # user-configured master
FIGDIR = ROOT / "figures"
FIGDIR.mkdir(exist_ok=True)

# ---- 1. Pre-render every EPS referenced (both sfigure/*.eps and figure/*.pstex) ----
def eps_to_png(eps: pathlib.Path, png: pathlib.Path):
    if png.exists(): return
    subprocess.run(
        ["gs", "-dQUIET", "-dNOPAUSE", "-dBATCH",
         "-sDEVICE=pngalpha", "-r200", "-dEPSCrop",
         f"-sOutputFile={png}", str(eps)],
        check=False,
    )

tex_raw = SRC.read_text(encoding="utf-8")

for m in re.finditer(r"\\includegraphics(?:\[[^\]]*\])?\{sfigure/([^}]+)\.eps\}", tex_raw):
    name = m.group(1)
    eps  = BASE.parent / "sfigure" / f"{name}.eps"
    if eps.exists():
        eps_to_png(eps, FIGDIR / f"{name}.png")

for m in re.finditer(r"\\input\{\.\./figure/([^}]+)\.pstex_t\}", tex_raw):
    name = m.group(1)
    eps  = BASE.parent / "figure" / f"{name}.pstex"
    if eps.exists():
        eps_to_png(eps, FIGDIR / f"{name}.png")

# placeholder blank png for tikz (user replaces later)
PLACE = FIGDIR / "_placeholder.png"
if not PLACE.exists():
    from PIL import Image, ImageDraw
    im = Image.new("RGB", (600, 300), "white")
    d = ImageDraw.Draw(im); d.rectangle((0,0,599,299), outline="gray", width=2)
    d.text((20,20), "(tikz placeholder)", fill="gray")
    im.save(PLACE)

# ---- 2. Preprocess tex ----
tex = tex_raw
tex = re.sub(r"\\Blue\{",  r"\\textcolor{blue}{", tex)
tex = re.sub(r"\\Red\{",   r"\\textcolor{red}{",  tex)

# --- custom math operators defined in custom.sty via \def\foo{\mathop{...}} ---
# pandoc doesn't evaluate these; substitute with \operatorname
for op in ("rank","Re","Im","det","adj","Span"):
    tex = re.sub(r"\\" + op + r"(?![A-Za-z])", r"\\operatorname{" + op.lower() + "}", tex)
# \defeq -> \stackrel{\text{def}}{=}
tex = re.sub(r"\\defeq(?![A-Za-z])", r"\\stackrel{\\text{def}}{=}", tex)

# pandoc's math parser accepts \nonumber but not \notag; normalize
tex = re.sub(r"\\notag(?![A-Za-z])", r"\\nonumber", tex)
# \cal -> \mathcal
tex = re.sub(r"\\\{cal\}", r"\\mathcal", tex)
tex = re.sub(r"\\cal(?![A-Za-z])\s*", r"\\mathcal ", tex)
# \\[5mm] etc. — strip the bracketed spacing in math line breaks
tex = re.sub(r"\\\\\s*\[[^\]]*\]", r"\\\\", tex)

# --- strip \hspace / \hspace* / \vspace inside equations (pandoc can't parse) ---
# we blanket-strip these since they're visual-only spacing
tex = re.sub(r"\\hspace\*?\{[^}]*\}", " ", tex)
tex = re.sub(r"\\vspace\*?\{[^}]*\}", " ", tex)

# figure pstex_t -> includegraphics
def fig_sub(m):
    name = pathlib.Path(m.group(1)).stem
    return r"\includegraphics[width=6cm]{figures/" + name + ".png}"
tex = re.sub(r"\\input\{([^}]+\.pstex_t)\}", fig_sub, tex)

# sfigure eps -> figures/NAME.png
tex = re.sub(
    r"\\includegraphics(\[[^\]]*\])?\{sfigure/([^}]+)\.eps\}",
    lambda m: r"\includegraphics" + (m.group(1) or "[width=6cm]") + "{figures/" + m.group(2) + ".png}",
    tex,
)

# tikzpicture blocks -> placeholder image (numbered)
tikz_idx = [0]
def tikz_sub(m):
    tikz_idx[0] += 1
    return r"\includegraphics[width=6cm]{figures/_placeholder.png}"
tex = re.sub(r"\\begin\{tikzpicture\}.*?\\end\{tikzpicture\}", tikz_sub, tex, flags=re.S)

# strip \scalebox{...}{ ... } wrapper (brace-aware)
def strip_scalebox(s):
    out=[]; i=0
    while i<len(s):
        if s.startswith(r"\scalebox{", i):
            j=i+len(r"\scalebox{"); d=1
            while j<len(s) and d:
                if s[j]=="{": d+=1
                elif s[j]=="}": d-=1
                j+=1
            while j<len(s) and s[j] in " \t\n": j+=1
            if j<len(s) and s[j]=="{":
                j+=1; d=1; st=j
                while j<len(s) and d:
                    if s[j]=="{": d+=1
                    elif s[j]=="}": d-=1
                    j+=1
                out.append(s[st:j-1]); i=j; continue
        out.append(s[i]); i+=1
    return "".join(out)
tex = strip_scalebox(tex)

# remove \tag{\ref{...}}
tex = re.sub(r"\\tag\{\\ref\{[^}]+\}\}", "", tex)

# strip \label{...} (cross-ref labels confuse pandoc's math OMML path)
tex = re.sub(r"\\label\{[^}]*\}", "", tex)
# \eqref{...} / \ref{...} -> plain text placeholder (not cross-referenced in pptx)
tex = re.sub(r"\\eqref\{[^}]*\}", "", tex)

# ---- 3. Build wrapped beamer tex ----
preamble = r"""
\documentclass[10pt]{beamer}
\usepackage[utf8]{inputenc}
\usepackage{amsmath,amssymb,bm,graphicx,xcolor}
\graphicspath{{./}}
\newcommand{\chap}[1]{}
\setcounter{chapter}{""" + str(N) + r"""}
\renewcommand{\theequation}{""" + str(N) + r""".\arabic{equation}}
\title{第""" + str(N) + r"""章}
\begin{document}
"""
postamble = r"\end{document}" + "\n"
WRAP.write_text(preamble + tex + "\n" + postamble, encoding="utf-8")

# ---- 4. Convert: tex -> md -> pptx (with slide-break normalization) ----
md = pypandoc.convert_file(str(WRAP), to="markdown", format="latex",
                           extra_args=[f"--resource-path={ROOT}"])
md = re.sub(r"^:::+.*$", "", md, flags=re.M)
md = re.sub(r"^### ", "## ", md, flags=re.M)
def _figrepl(m):
    img = re.search(r'<img\s+src="([^"]+)"', m.group(0))
    cap = re.search(r'<figcaption>(.*?)</figcaption>', m.group(0), re.S)
    if not img: return ""
    return f"![{cap.group(1) if cap else ''}]({img.group(1)})"
md = re.sub(r"<figure[^>]*>.*?</figure>", _figrepl, md, flags=re.S)
MD.write_text(md, encoding="utf-8")

extra = ["--slide-level=2", f"--resource-path={ROOT}"]
if REF.exists() and N != 3:
    extra += [f"--reference-doc={REF}"]
pypandoc.convert_file(str(MD), to="pptx", format="markdown",
                      outputfile=str(OUT), extra_args=extra)
print(f"wrote: {OUT}")

# ---- 5. Post-process: namespace fix + autofit + picture layout ----
import zipfile, shutil
from lxml import etree
_NS_A  = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_NS_P  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
_QA = '{%s}' % _NS_A
_QP = '{%s}' % _NS_P

def _flatten_cwc(xml_bytes: bytes) -> bytes:
    """Merge Content-with-Caption residue (idx=2 sz=half body + idx=1 textbox
    + idx=1 content) into a single Title-and-Content body placeholder."""
    tree = etree.fromstring(xml_bytes)
    spTree = tree.find('.//' + _QP + 'cSld/' + _QP + 'spTree')
    if spTree is None: return xml_bytes
    caption = textbox = content = None
    for sp in spTree.iter(_QP + 'sp'):
        ph = sp.find('.//' + _QP + 'nvSpPr/' + _QP + 'nvPr/' + _QP + 'ph')
        if ph is None: continue
        idx, sz, typ = ph.get('idx'), ph.get('sz'), ph.get('type')
        is_txbox = (sp.find(_QP + 'nvSpPr/' + _QP + 'cNvSpPr') or etree.Element('x')).get('txBox') == '1'
        if idx == '2' and sz == 'half' and typ == 'body':
            caption = sp
        elif idx == '1' and is_txbox:
            textbox = sp
        elif idx == '1':
            content = sp
    if caption is None: return xml_bytes
    paras = []
    for sp in (caption, textbox, content):
        if sp is None: continue
        txb = sp.find(_QP + 'txBody')
        if txb is None: continue
        paras.extend(txb.findall(_QA + 'p'))
    if content is None:
        content = caption
        ph = content.find('.//' + _QP + 'nvSpPr/' + _QP + 'nvPr/' + _QP + 'ph')
        for k in ('sz', 'type'):
            ph.attrib.pop(k, None)
        ph.set('idx', '1')
        cNvPr = content.find(_QP + 'nvSpPr/' + _QP + 'cNvPr')
        if cNvPr is not None: cNvPr.set('name', 'Content Placeholder 2')
    target = content.find(_QP + 'txBody')
    for p in target.findall(_QA + 'p'):
        target.remove(p)
    for p in paras:
        target.append(p)
    for sp in (caption, textbox):
        if sp is None or sp is content: continue
        node = sp
        while node.getparent() is not spTree:
            node = node.getparent()
            if node is None: break
        if node is not None and node.getparent() is spTree:
            spTree.remove(node)
    return etree.tostring(tree, xml_declaration=True, encoding='UTF-8', standalone=True)

A14 = 'xmlns:a14="http://schemas.microsoft.com/office/drawing/2010/main" '
MC  = 'xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
FONT_SCALE = 62000
LINE_REDUCTION = 20000
AUTOFIT = f'<a:normAutofit fontScale="{FONT_SCALE}" lnSpcReduction="{LINE_REDUCTION}"/>'

# Identify the rId# of "Title and Content" layout from the master's rels
_TAC_LAYOUT = None
with zipfile.ZipFile(OUT,"r") as _z:
    for _n in _z.namelist():
        if _n.startswith("ppt/slideLayouts/") and _n.endswith(".xml"):
            _x = _z.read(_n).decode("utf-8")
            if '<p:cSld name="Title and Content"' in _x:
                _TAC_LAYOUT = _n.split("/")[-1]; break
print(f"Title-and-Content layout = {_TAC_LAYOUT}")

tmp = OUT.with_suffix(".pptx.tmp")
with zipfile.ZipFile(OUT,"r") as zi, zipfile.ZipFile(tmp,"w",zipfile.ZIP_DEFLATED) as zo:
    for it in zi.infolist():
        data = zi.read(it.filename)
        # For every slide (except slide1 = title page), force layout -> Title and Content
        m_rel = re.match(r"ppt/slides/_rels/slide(\d+)\.xml\.rels$", it.filename)
        if m_rel and _TAC_LAYOUT:
            s = data.decode("utf-8")
            s = re.sub(
                r'(Type="http://schemas\.openxmlformats\.org/officeDocument/2006/relationships/slideLayout"\s+Target=")[^"]+(")',
                r'\1../slideLayouts/' + _TAC_LAYOUT + r'\2',
                s, count=1,
            )
            data = s.encode("utf-8")
        if re.match(r"ppt/slides/slide\d+\.xml$", it.filename):
            s = data.decode("utf-8")
            # ensure a14/mc namespaces declared on root so math in tables is valid
            def _addns(m):
                a = m.group(1); add = ""
                if "xmlns:a14=" not in a: add += " " + A14
                if "xmlns:mc="  not in a: add += " " + MC
                return "<p:sld" + a + add + ">"
            s = re.sub(r"<p:sld([^>]*)>", _addns, s, count=1)
            # aggressive shrink-to-fit on every bodyPr
            s = re.sub(r"<a:normAutofit[^/]*/>", AUTOFIT, s)
            s = re.sub(r"<a:bodyPr([^/]*)/>",
                       lambda m: f"<a:bodyPr{m.group(1)}>{AUTOFIT}</a:bodyPr>"
                       if "normAutofit" not in m.group(1) else m.group(0),
                       s)
            data = s.encode("utf-8")
            # flatten Content-with-Caption residue into single Title+Content body
            data = _flatten_cwc(data)
        zo.writestr(it, data)
shutil.move(tmp, OUT)

# fit pictures into the layout's content box (text-left / picture-right)
from pptx import Presentation as _P
from pptx.util import Emu as _Emu
p = _P(OUT)
def _cbox(sl):
    for ph in sl.slide_layout.placeholders:
        if ph.placeholder_format.idx == 1:
            return ph.left, ph.top, ph.width, ph.height
    return _Emu(133304), _Emu(685800), _Emu(8959895), _Emu(4019549)
GAP = _Emu(91440)
for sl in p.slides:
    cx,cy,cw,ch = _cbox(sl)
    pics = [s for s in sl.shapes if s.shape_type == 13]
    for pic in pics:
        r = pic.width / pic.height if pic.height else 1.0
        tw = (cw - GAP)//2; pw = cw - GAP - tw
        if pw/r <= ch: nw, nh = pw, int(pw/r)
        else:           nh, nw = ch, int(ch*r)
        pic.width, pic.height = nw, nh
        pic.left = cx + tw + GAP + (pw - nw)//2
        pic.top  = cy + (ch - nh)//2
p.save(OUT)
print(f"post-processed: {OUT}")
