#!/usr/bin/env python3
"""Post-process chap3.pptx so all content fits inside the master layout area.

Strategy:
  - For every body placeholder (incl. those pandoc wraps in
    <mc:AlternateContent>), tighten <a:normAutofit> to a stronger shrink
    so overflowing math/text is scaled down to fit the content box.
  - For slides with a picture alongside text, split the layout content area
    horizontally (text left, picture right) and resize the picture to fit.
"""
import re, zipfile, shutil, pathlib
from pptx import Presentation
from pptx.util import Emu

HERE = pathlib.Path(__file__).parent
PPTX = HERE / "chap3.pptx"

# ---- 1. Patch normAutofit in every slide XML (handles AlternateContent wrap) ----
FONT_SCALE = 62000        # 62%
LINE_REDUCTION = 20000    # 20% line spacing cut
AUTOFIT_ATTRS = f' fontScale="{FONT_SCALE}" lnSpcReduction="{LINE_REDUCTION}"'

tmp = PPTX.with_suffix(".pptx.tmp")
with zipfile.ZipFile(PPTX, "r") as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
    for item in zin.infolist():
        data = zin.read(item.filename)
        if re.match(r"ppt/slides/slide\d+\.xml$", item.filename):
            s = data.decode("utf-8")
            # Normalize every <a:normAutofit .../> to our stronger shrink.
            s = re.sub(r"<a:normAutofit[^/]*/>", f"<a:normAutofit{AUTOFIT_ATTRS}/>", s)
            # Some bodyPr elements have no normAutofit; inject one.
            def _inject(m):
                inner = m.group(1)
                if "normAutofit" in inner:
                    return m.group(0)
                return f"<a:bodyPr{inner}><a:normAutofit{AUTOFIT_ATTRS}/></a:bodyPr>"
            s = re.sub(r"<a:bodyPr([^/]*)/>", _inject, s)
            data = s.encode("utf-8")
        zout.writestr(item, data)
shutil.move(tmp, PPTX)
print(f"patched autofit in {PPTX.name}")

# ---- 2. Resize / reposition pictures using python-pptx ----
p = Presentation(PPTX)

def layout_content_box(slide):
    for ph in slide.slide_layout.placeholders:
        if ph.placeholder_format.idx == 1:
            return ph.left, ph.top, ph.width, ph.height
    return Emu(133304), Emu(685800), Emu(8959895), Emu(4019549)

GAP = Emu(91440)  # 0.1"

for slide in p.slides:
    cx, cy, cw, ch = layout_content_box(slide)
    pictures = [s for s in slide.shapes if s.shape_type == 13]
    if not pictures:
        continue

    # Detect whether this slide has body text (via AlternateContent the shape
    # may not be visible to python-pptx — check raw XML for <a:t> outside title)
    # Heuristic: if picture's existing left > cw/3, assume two-column layout.
    pic = pictures[0]
    has_text_left = True

    ratio = pic.width / pic.height if pic.height else 1.0
    if has_text_left:
        text_w = (cw - GAP) // 2
        pic_w  = cw - GAP - text_w
        max_w, max_h = pic_w, ch
        if max_w / ratio <= max_h:
            new_w, new_h = max_w, int(max_w / ratio)
        else:
            new_h, new_w = max_h, int(max_h * ratio)
        pic.width, pic.height = new_w, new_h
        pic.left = cx + text_w + GAP + (pic_w - new_w) // 2
        pic.top  = cy + (ch - new_h) // 2
    else:
        if cw / ratio <= ch:
            new_w, new_h = cw, int(cw / ratio)
        else:
            new_h, new_w = ch, int(ch * ratio)
        pic.width, pic.height = new_w, new_h
        pic.left = cx + (cw - new_w) // 2
        pic.top  = cy + (ch - new_h) // 2

p.save(PPTX)
print(f"repositioned pictures in {PPTX.name}")
